#!/usr/bin/env python3
"""Debug PowerPoint generation - trace section processing."""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent / "src"))

from md2office.parser import MarkdownParser, ASTBuilder
from md2office.generators.powerpoint_generator import PowerPointGenerator
from md2office.parser.ast_builder import NodeType

md_file = Path("examples/input/formation-copilot-365.md")
content = md_file.read_text(encoding='utf-8')

parser = MarkdownParser()
tokens = parser.parse(content)
builder = ASTBuilder()
ast = builder.build(tokens)

generator = PowerPointGenerator()

# Monkey patch _process_section to add logging
original_process_section = generator._process_section

def logged_process_section(node, options):
    level = node.level or 1
    heading = None
    for child in node.children:
        if child.node_type == NodeType.HEADING:
            heading = child
            break
    
    heading_text = heading.content if heading else "No heading"
    print(f"Processing SECTION level={level}: {heading_text}")
    
    result = original_process_section(node, options)
    
    print(f"  After processing level={level}, current_slide exists: {generator.current_slide is not None}")
    if generator.current_slide:
        try:
            title = generator.current_slide.shapes.title.text if generator.current_slide.shapes.title else "No title"
            print(f"  Current slide title: {title}")
        except:
            pass
    
    return result

generator._process_section = logged_process_section

# Generate
print("Starting PowerPoint generation...")
print("=" * 60)
pptx_bytes = generator.generate(ast, {})

from pptx import Presentation
from io import BytesIO
pres = Presentation(BytesIO(pptx_bytes))
print(f"\nFinal slide count: {len(pres.slides)}")

