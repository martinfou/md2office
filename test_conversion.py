#!/usr/bin/env python3
"""
Test script to verify PowerPoint conversion works correctly.
"""
import sys
from pathlib import Path

# Add src to path
sys.path.insert(0, str(Path(__file__).parent / "src"))

try:
    from md2office.parser import MarkdownParser, ASTBuilder
    from md2office.generators import PowerPointGenerator
    
    print("Testing PowerPoint conversion...")
    print("=" * 50)
    
    # Read the markdown file
    md_file = Path("examples/input/formation-copilot-365.md")
    if not md_file.exists():
        print(f"ERROR: Markdown file not found: {md_file}")
        sys.exit(1)
    
    print(f"Reading markdown file: {md_file}")
    content = md_file.read_text(encoding='utf-8')
    
    # Parse markdown
    print("Parsing markdown...")
    parser = MarkdownParser()
    tokens = parser.parse(content)
    print(f"  Parsed {len(tokens)} tokens")
    
    # Build AST
    print("Building AST...")
    builder = ASTBuilder()
    ast = builder.build(tokens)
    print(f"  AST built with {len(ast.children)} top-level children")
    
    # Count sections
    def count_sections(node):
        count = 0
        if node.node_type.value == "section":
            count = 1
        for child in node.children:
            count += count_sections(child)
        return count
    
    section_count = count_sections(ast)
    print(f"  Found {section_count} sections")
    
    # Generate PowerPoint
    print("Generating PowerPoint...")
    generator = PowerPointGenerator()
    pptx_bytes = generator.generate(ast, {})
    print(f"  Generated {len(pptx_bytes)} bytes")
    
    # Save to file
    output_file = Path("formation-copilot-365-test.pptx")
    output_file.write_bytes(pptx_bytes)
    print(f"  Saved to: {output_file}")
    
    # Verify it's a valid PPTX file (starts with ZIP signature)
    if pptx_bytes[:2] == b'PK':
        print("  ✓ Valid PPTX file (ZIP format)")
    else:
        print("  ✗ Invalid PPTX file format")
        sys.exit(1)
    
    # Try to read the presentation to count slides
    try:
        from pptx import Presentation
        from io import BytesIO
        pres = Presentation(BytesIO(pptx_bytes))
        slide_count = len(pres.slides)
        print(f"  ✓ Presentation has {slide_count} slides")
        
        if slide_count == 1:
            print("\n⚠️  WARNING: Only 1 slide generated! This suggests the fix may not be working.")
            print("   Expected multiple slides (one per H2 section)")
        elif slide_count > 1:
            print(f"\n✓ SUCCESS: Generated {slide_count} slides correctly!")
        else:
            print("\n✗ ERROR: No slides generated!")
            sys.exit(1)
            
        # Print slide titles
        print("\nSlide titles:")
        for i, slide in enumerate(pres.slides, 1):
            title = "No title"
            if slide.shapes.title:
                title = slide.shapes.title.text
            print(f"  Slide {i}: {title}")
            
    except ImportError:
        print("  (Cannot verify slide count - python-pptx not available for reading)")
    
    print("\n" + "=" * 50)
    print("Conversion test completed successfully!")
    
except ImportError as e:
    print(f"ERROR: Missing dependency: {e}")
    print("\nPlease install dependencies:")
    print("  pip install python-pptx")
    sys.exit(1)
except Exception as e:
    print(f"ERROR: {e}")
    import traceback
    traceback.print_exc()
    sys.exit(1)

