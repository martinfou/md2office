"""
PowerPoint Document Generator

Implements Epic 3: PowerPoint Conversion
Generates Microsoft PowerPoint (.pptx) presentations from AST.
"""

from typing import Dict, Any, Optional, List, Tuple
from io import BytesIO
from pathlib import Path
import re
import subprocess
import tempfile
import hashlib

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image as PILImage
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

from ..parser.ast_builder import ASTNode, NodeType
from ..router.content_router import FormatGenerator, OutputFormat
from ..errors import ConversionError, FileError
from ..styling.style import StylePreset, get_style_preset


class PowerPointGenerator(FormatGenerator):
    """
    Generates PowerPoint (.pptx) presentations from AST.
    
    Implements Story 3.1: PowerPoint Document Generator Core
    """
    
    def __init__(self):
        """Initialize PowerPoint generator."""
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PowerPoint generation. "
                "Install with: pip install python-pptx"
            )
        self.presentation: Optional[Presentation] = None
        self.current_slide = None
        self.current_style_preset: Optional[StylePreset] = None
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
        self._temp_dir: Optional[tempfile.TemporaryDirectory] = None  # Keep temp dir alive during generation
        
        # Patterns for inline markdown parsing
        self.bold_pattern = re.compile(r'\*\*([^*]+)\*\*|__([^_]+)__')
        self.italic_pattern = re.compile(r'(?<!\*)\*([^*]+?)\*(?!\*)|(?<!_)_([^_]+?)_(?!_)')
        self.code_pattern = re.compile(r'`([^`]+)`')
        self.link_pattern = re.compile(r'\[([^\]]+)\]\(([^)]+)\)')
    
    def generate(self, ast: ASTNode, options: Dict[str, Any]) -> bytes:
        """
        Generate PowerPoint presentation from AST.
        
        Args:
            ast: Root AST node
            options: Generation options
            
        Returns:
            Generated PowerPoint presentation as bytes
            
        Raises:
            ConversionError: If generation fails
        """
        # Create temporary directory for Mermaid images (kept alive during generation)
        self._temp_dir = tempfile.TemporaryDirectory()
        
        try:
            # Initialize presentation
            self.presentation = Presentation()
            self.presentation.slide_width = self.slide_width
            self.presentation.slide_height = self.slide_height
            
            # Get style preset
            style_name = options.get('style', 'default')
            self.current_style_preset = get_style_preset(style_name)
            
            # Set presentation metadata
            self._set_presentation_metadata(ast, options)
            
            # Process AST nodes
            self._process_node(ast, options)
            
            # Save to bytes
            output = BytesIO()
            self.presentation.save(output)
            return output.getvalue()
        
        except Exception as e:
            raise ConversionError(
                f"Failed to generate PowerPoint presentation: {str(e)}",
                format="powerpoint",
                stage="generation"
            ) from e
        finally:
            # Clean up temporary directory
            if self._temp_dir:
                self._temp_dir.cleanup()
                self._temp_dir = None
    
    def get_file_extension(self) -> str:
        """Get file extension for PowerPoint format."""
        return ".pptx"
    
    def _set_presentation_metadata(self, ast: ASTNode, options: Dict[str, Any]):
        """Set presentation metadata."""
        # Get title from AST metadata or first H1
        title = ast.metadata.get('title')
        if not title:
            first_h1 = self._find_first_heading(ast, 1)
            if first_h1:
                title = first_h1.content
        
        if title:
            self.presentation.core_properties.title = title
        
        # Get author from metadata or options
        author = ast.metadata.get('author') or options.get('author')
        if author:
            self.presentation.core_properties.author = author
    
    def _process_node(self, node: ASTNode, options: Dict[str, Any]):
        """
        Process AST node and add to presentation.
        
        Args:
            node: AST node to process
            options: Generation options
        """
        if node.node_type == NodeType.DOCUMENT:
            # Process all children of document
            for child in node.children:
                self._process_node(child, options)
        
        elif node.node_type == NodeType.SECTION:
            # Process section - may create new slide
            # Section processing handles its own children, so don't recurse here
            self._process_section(node, options)
        
        elif node.node_type == NodeType.HEADING:
            self._process_heading(node, options)
        
        elif node.node_type == NodeType.PARAGRAPH:
            self._add_paragraph_to_slide(node, options)
        
        elif node.node_type == NodeType.LIST:
            self._add_list_to_slide(node, options)
        
        elif node.node_type == NodeType.TABLE:
            self._add_table_to_slide(node, options)
        
        elif node.node_type == NodeType.CODE_BLOCK:
            self._add_code_block_to_slide(node, options)
        
        elif node.node_type == NodeType.BLOCKQUOTE:
            self._add_blockquote_to_slide(node, options)
        
        elif node.node_type == NodeType.IMAGE:
            self._add_image_to_slide(node, options)
        
        # Process children recursively (but not for SECTION nodes - they handle their own children)
        # Also skip already-processed content nodes
        for child in node.children:
            if child.node_type not in [NodeType.HEADING, NodeType.PARAGRAPH,
                                      NodeType.LIST, NodeType.TABLE,
                                      NodeType.CODE_BLOCK, NodeType.BLOCKQUOTE,
                                      NodeType.IMAGE, NodeType.SECTION]:
                self._process_node(child, options)
    
    def _process_section(self, node: ASTNode, options: Dict[str, Any]):
        """Process section node - creates slides based on heading level."""
        level = node.level or 1
        
        # Find heading in section
        heading = None
        for child in node.children:
            if child.node_type == NodeType.HEADING:
                heading = child
                break
        
        if level == 1:
            # H1 section - create title slide
            self._create_title_slide(heading, options)
        elif level == 2:
            # H2 section - create section header slide
            self._create_section_slide(heading, options)
        else:
            # H3+ section - may create content slide or add to current slide
            self._create_content_slide(heading, options)
        
        # Process section content (excluding heading and nested sections)
        # Nested sections will be processed separately to create their own slides
        # Note: Sections can be nested under headings, so we need to check heading children too
        for child in node.children:
            if child.node_type == NodeType.HEADING:
                # Heading already processed above, but check if it has nested sections as children
                for heading_child in child.children:
                    if heading_child.node_type == NodeType.SECTION:
                        # Process nested sections that are children of the heading
                        self._process_section(heading_child, options)
                    else:
                        # Process other heading children (content)
                        self._process_node(heading_child, options)
            elif child.node_type == NodeType.SECTION:
                # Process nested sections to create new slides
                self._process_section(child, options)
            else:
                # Process content nodes (paragraphs, lists, etc.)
                self._process_node(child, options)
    
    def _process_heading(self, node: ASTNode, options: Dict[str, Any]):
        """Process heading node."""
        level = node.level or 1
        
        if level == 1:
            self._create_title_slide(node, options)
        elif level == 2:
            self._create_section_slide(node, options)
        else:
            # H3+ - add as slide title or content header
            if not self.current_slide:
                self._create_content_slide(node, options)
            else:
                # Add as subsection header in current slide
                self._add_subsection_header(node, options)
    
    def _create_title_slide(self, heading: Optional[ASTNode], options: Dict[str, Any]):
        """Create title slide from H1 heading."""
        layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(layout)
        self.current_slide = slide
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        if heading:
            title.text = heading.content
        else:
            title.text = "Presentation"
        
        # Set subtitle if available
        if subtitle:
            # Try to get subtitle from metadata or first paragraph
            subtitle_text = options.get('subtitle', '')
            if subtitle_text:
                subtitle.text = subtitle_text
    
    def _create_section_slide(self, heading: Optional[ASTNode], options: Dict[str, Any]):
        """Create section header slide from H2 heading."""
        layout = self.presentation.slide_layouts[1]  # Title and Content layout
        slide = self.presentation.slides.add_slide(layout)
        self.current_slide = slide
        
        title = slide.shapes.title
        if heading:
            title.text = heading.content
        else:
            title.text = "Section"
    
    def _create_content_slide(self, heading: Optional[ASTNode], options: Dict[str, Any]):
        """Create content slide."""
        layout = self.presentation.slide_layouts[1]  # Title and Content layout
        slide = self.presentation.slides.add_slide(layout)
        self.current_slide = slide
        
        title = slide.shapes.title
        if heading:
            title.text = heading.content
        else:
            title.text = ""
    
    def _add_subsection_header(self, node: ASTNode, options: Dict[str, Any]):
        """Add subsection header to current slide."""
        if not self.current_slide:
            return
        
        # Find content placeholder
        content_placeholder = None
        for shape in self.current_slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            text_frame = content_placeholder.text_frame
            p = text_frame.add_paragraph()
            p.level = (node.level or 3) - 2  # Convert H3+ to paragraph level
            # Parse and apply inline formatting (subsection headers are bold by default)
            self._add_formatted_text(p, node.content, base_font_size=Pt(20), default_bold=True)
    
    def _slide_has_code_or_image(self) -> bool:
        """Check if current slide has code blocks or images (non-placeholder shapes)."""
        if not self.current_slide:
            return False
        
        # Check for shapes that are not placeholders (code blocks, images)
        for shape in self.current_slide.shapes:
            try:
                # Skip title placeholder (idx 0) and content placeholder (idx 1)
                if hasattr(shape, 'placeholder_format'):
                    try:
                        _ = shape.placeholder_format.idx  # Test if it's actually a placeholder
                        continue
                    except:
                        # Not actually a placeholder, continue checking
                        pass
            except:
                pass
            
            # Check if it's an image
            try:
                if hasattr(shape, 'image'):
                    return True
            except:
                pass
            
            # Check if it's a textbox with code-like background (gray background indicates code block)
            try:
                if hasattr(shape, 'fill') and hasattr(shape, 'text_frame'):
                    if hasattr(shape.fill, 'fore_color'):
                        # Check if it has a light gray background (code blocks use #F5F5F5)
                        if hasattr(shape.fill.fore_color, 'rgb'):
                            rgb = shape.fill.fore_color.rgb
                            # RGBColor can be accessed as tuple (r, g, b) or has .r, .g, .b attributes
                            try:
                                if rgb:
                                    # Try to access as tuple/index
                                    if hasattr(rgb, '__len__') and len(rgb) >= 3:
                                        r, g, b = rgb[0], rgb[1], rgb[2]
                                    elif hasattr(rgb, 'r') and hasattr(rgb, 'g') and hasattr(rgb, 'b'):
                                        r, g, b = rgb.r, rgb.g, rgb.b
                                    else:
                                        continue
                                    # Check if it's approximately gray (#F5F5F5 = 245, 245, 245)
                                    if r >= 240 and g >= 240 and b >= 240:
                                        return True
                            except (TypeError, AttributeError, IndexError):
                                pass
            except:
                pass
        
        return False
    
    def _reorganize_slide_for_text_after_code_image(self):
        """Reorganize slide when text is added after code block or image: move code/image to right, text to left."""
        if not self.current_slide:
            return
        
        # Find content placeholder and move to left column
        content_placeholder = None
        for shape in self.current_slide.shapes:
            try:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:
                    content_placeholder = shape
                    break
            except:
                continue
        
        if content_placeholder:
            try:
                # Move content placeholder to left column
                content_placeholder.left = Inches(0.5)
                content_placeholder.top = Inches(2)
                content_placeholder.width = Inches(4)
                content_placeholder.height = Inches(5)
                
                # Enable word wrap
                if hasattr(content_placeholder, 'text_frame'):
                    content_placeholder.text_frame.word_wrap = True
            except:
                pass
        
        # Move existing code blocks and images to right column
        for shape in self.current_slide.shapes:
            # Skip placeholders - check carefully to avoid exceptions
            is_placeholder = False
            try:
                if hasattr(shape, 'placeholder_format'):
                    try:
                        _ = shape.placeholder_format.idx  # Test if it's actually a placeholder
                        is_placeholder = True
                    except (AttributeError, ValueError):
                        # Not actually a placeholder, continue checking
                        pass
            except:
                # If we can't check, assume it's not a placeholder and continue
                pass
            
            if is_placeholder:
                continue
            
            # Move code blocks and images to right column
            is_code_or_image = False
            
            # Check if it's an image
            try:
                if hasattr(shape, 'image'):
                    is_code_or_image = True
            except:
                pass
            
            # Check if it's a code block (textbox with gray background)
            if not is_code_or_image:
                try:
                    if hasattr(shape, 'fill') and hasattr(shape, 'text_frame'):
                        # Try to access fill color - code blocks have solid fill with gray background
                        try:
                            if hasattr(shape.fill, 'fore_color'):
                                try:
                                    if hasattr(shape.fill.fore_color, 'rgb'):
                                        rgb = shape.fill.fore_color.rgb
                                        if rgb:
                                            # RGBColor can be accessed as tuple (r, g, b) or has .r, .g, .b attributes
                                            try:
                                                if hasattr(rgb, '__len__') and len(rgb) >= 3:
                                                    r, g, b = rgb[0], rgb[1], rgb[2]
                                                elif hasattr(rgb, 'r') and hasattr(rgb, 'g') and hasattr(rgb, 'b'):
                                                    r, g, b = rgb.r, rgb.g, rgb.b
                                                else:
                                                    # Can't determine RGB values, skip
                                                    pass
                                                # Check if it's approximately gray (#F5F5F5 = 245, 245, 245)
                                                if r >= 240 and g >= 240 and b >= 240:
                                                    is_code_or_image = True
                                            except (TypeError, AttributeError, IndexError):
                                                pass
                                except (AttributeError, ValueError):
                                    # Fill doesn't have accessible fore_color (e.g., NoneFill) - not a code block
                                    pass
                        except (AttributeError, ValueError):
                            # Can't access fill properties
                            pass
                except:
                    pass
            
            if is_code_or_image:
                # Move to right column
                try:
                    shape.left = Inches(5.5)
                    shape.top = Inches(2)
                    # Adjust width to fit right column
                    if hasattr(shape, 'width') and shape.width > Inches(4):
                        shape.width = Inches(4)
                    if hasattr(shape, 'height') and shape.height > Inches(5):
                        shape.height = Inches(5)
                except Exception as e:
                    # If moving fails, skip this shape (might be locked or have restrictions)
                    pass
    
    def _add_paragraph_to_slide(self, node: ASTNode, options: Dict[str, Any]):
        """Add paragraph to current slide."""
        if not self.current_slide:
            self._create_content_slide(None, options)
        
        # Check if there are already code blocks or images on the slide
        if self._slide_has_code_or_image():
            self._reorganize_slide_for_text_after_code_image()
        
        # Find content placeholder
        content_placeholder = None
        for shape in self.current_slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            text_frame = content_placeholder.text_frame
            p = text_frame.add_paragraph()
            p.space_after = Pt(12)
            # Parse and apply inline formatting
            self._add_formatted_text(p, node.content, base_font_size=Pt(18))
    
    def _add_list_to_slide(self, node: ASTNode, options: Dict[str, Any]):
        """Add list to current slide."""
        if not self.current_slide:
            self._create_content_slide(None, options)
        
        # Check if there are already code blocks or images on the slide
        if self._slide_has_code_or_image():
            self._reorganize_slide_for_text_after_code_image()
        
        # Find content placeholder
        content_placeholder = None
        for shape in self.current_slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            text_frame = content_placeholder.text_frame
            
            for list_item in node.children:
                if list_item.node_type == NodeType.LIST_ITEM:
                    p = text_frame.add_paragraph()
                    p.level = list_item.level or 0
                    
                    # Set bullet style
                    is_ordered = list_item.metadata.get('ordered', False)
                    content = list_item.content
                    if is_ordered:
                        # Numbered list (simplified - full implementation would handle numbering)
                        marker = list_item.metadata.get('marker', '1.')
                        # Parse and apply inline formatting
                        self._add_formatted_text(p, f"{marker} {content}", base_font_size=Pt(18))
                    else:
                        # Parse and apply inline formatting
                        self._add_formatted_text(p, content, base_font_size=Pt(18))
    
    def _add_table_to_slide(self, node: ASTNode, options: Dict[str, Any]):
        """Add table to current slide."""
        if not self.current_slide:
            self._create_content_slide(None, options)
        
        if not node.metadata:
            return
        
        headers = node.metadata.get('headers', [])
        rows = node.metadata.get('rows', [])
        
        if not headers:
            return
        
        # Create table shape
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(min(4, len(rows) * 0.5 + 1))
        
        table_shape = self.current_slide.shapes.add_table(
            rows=len(rows) + 1,
            cols=len(headers),
            left=left,
            top=top,
            width=width,
            height=height
        )
        
        table = table_shape.table
        
        # Add header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
            
            # Make header text bold
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
        
        # Add data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                if col_idx < len(headers):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = cell_data
    
    def _add_code_block_to_slide(self, node: ASTNode, options: Dict[str, Any]):
        """Add code block to current slide. If text exists, create two-column layout."""
        if not self.current_slide:
            self._create_content_slide(None, options)
        
        content = node.content
        # Handle case where language might be None
        language = (node.metadata.get('language') or '').lower()
        
        # Check if this is a Mermaid diagram
        is_mermaid = language == 'mermaid' or self._is_mermaid_diagram(content)
        
        if is_mermaid:
            # Try to render Mermaid diagram as image
            image_path = self._render_mermaid_diagram(content, options)
            if image_path and image_path.exists():
                # Add as image
                self._add_mermaid_image_to_slide(image_path, options)
                return
        
        # Check if there's existing text content on the slide
        has_text_content = self._slide_has_text_content()
        
        if has_text_content:
            # Create two-column layout: text on left, code on right
            self._reorganize_slide_for_two_columns()
            
            # Create text box for code in right column
            left = Inches(5.5)
            top = Inches(2)
            width = Inches(4)
            height = Inches(5)
        else:
            # No text content, use full width
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
        
        text_box = self.current_slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Set background color
        fill = text_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        
        # Add code text
        p = text_frame.paragraphs[0]
        p.text = content
        p.font.name = 'Courier New'
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    
    def _is_mermaid_diagram(self, content: str) -> bool:
        """Check if content looks like a Mermaid diagram."""
        if not content:
            return False
        
        first_line = content.strip().split('\n')[0].strip()
        mermaid_keywords = [
            'graph', 'flowchart', 'sequenceDiagram', 'classDiagram',
            'stateDiagram', 'stateDiagram-v2', 'erDiagram', 'gantt',
            'pie', 'gitgraph', 'journey', 'mindmap', 'C4Context',
            'quadrantChart', 'requirement', 'timeline'
        ]
        return any(first_line.startswith(keyword) for keyword in mermaid_keywords)
    
    def _render_mermaid_diagram(self, mermaid_code: str, options: Dict[str, Any]) -> Optional[Path]:
        """
        Render Mermaid diagram to image file.
        
        Returns:
            Path to rendered image file, or None if rendering failed
        """
        # Try multiple rendering methods
        # Method 1: Try mermaid-cli (mmdc) if available
        image_path = self._render_with_mermaid_cli(mermaid_code, options)
        if image_path:
            return image_path
        
        # Method 2: Try Playwright if available
        image_path = self._render_with_playwright(mermaid_code, options)
        if image_path:
            return image_path
        
        # If all methods fail, return None (will fall back to code display)
        return None
    
    def _render_with_mermaid_cli(self, mermaid_code: str, options: Dict[str, Any]) -> Optional[Path]:
        """Render Mermaid diagram using mermaid-cli (mmdc)."""
        try:
            # Check if mmdc is available
            result = subprocess.run(
                ['mmdc', '--version'],
                capture_output=True,
                timeout=5
            )
            if result.returncode != 0:
                return None
        except (FileNotFoundError, subprocess.TimeoutExpired):
            return None
        
        if not self._temp_dir:
            return None
        
        # Use generator's temp directory
        temp_dir_path = Path(self._temp_dir.name)
        
        # Create unique filename for this diagram
        diagram_hash = hashlib.md5(mermaid_code.encode()).hexdigest()[:8]
        input_file = temp_dir_path / f'diagram_{diagram_hash}.mmd'
        output_file = temp_dir_path / f'diagram_{diagram_hash}.png'
        
        # Write Mermaid code to file
        input_file.write_text(mermaid_code, encoding='utf-8')
        
        # Render diagram
        try:
            result = subprocess.run(
                [
                    'mmdc',
                    '-i', str(input_file),
                    '-o', str(output_file),
                    '-w', '1200',
                    '-H', '800',
                    '-b', 'transparent'
                ],
                capture_output=True,
                timeout=30,
                check=True
            )
            
            if output_file.exists():
                return output_file
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
            return None
        
        return None
    
    def _render_with_playwright(self, mermaid_code: str, options: Dict[str, Any]) -> Optional[Path]:
        """Render Mermaid diagram using Playwright."""
        try:
            from playwright.sync_api import sync_playwright
        except ImportError:
            return None
        
        if not self._temp_dir:
            return None
        
        try:
            temp_dir_path = Path(self._temp_dir.name)
            
            # Create unique filename for this diagram
            diagram_hash = hashlib.md5(mermaid_code.encode()).hexdigest()[:8]
            output_file = temp_dir_path / f'diagram_{diagram_hash}.png'
            
            html_content = f"""<!DOCTYPE html>
<html>
<head>
    <script src="https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.min.js"></script>
</head>
<body>
    <div class="mermaid">
{mermaid_code}
    </div>
    <script>
        mermaid.initialize({{ startOnLoad: true }});
    </script>
</body>
</html>"""
            
            with sync_playwright() as p:
                browser = p.chromium.launch()
                page = browser.new_page()
                page.set_content(html_content)
                page.wait_for_selector('.mermaid svg', timeout=10000)
                page.screenshot(path=str(output_file), full_page=True)
                browser.close()
            
            if output_file.exists():
                return output_file
        except Exception:
            return None
        
        return None
    
    def _add_mermaid_image_to_slide(self, image_path: Path, options: Dict[str, Any]):
        """Add rendered Mermaid diagram image to current slide. If text exists, create two-column layout."""
        if not self.current_slide:
            return
        
        try:
            # Check if there's existing text content on the slide
            has_text_content = self._slide_has_text_content()
            
            if has_text_content:
                # Create two-column layout: text on left, image on right
                self._reorganize_slide_for_two_columns()
                
                # Calculate dimensions preserving aspect ratio for right column
                max_width = 4.0  # Right column width in inches
                max_height = 5.5  # Available height in inches
                width, height = self._calculate_image_dimensions(image_path, max_width, max_height)
                
                # Position in right column (centered vertically)
                left = Inches(5.5)
                top = Inches(1.5 + (max_height - height) / 2)  # Center vertically
            else:
                # No text content, center the diagram
                max_width = 9.0  # Full width minus margins
                max_height = 5.5  # Available height
                width, height = self._calculate_image_dimensions(image_path, max_width, max_height)
                
                # Center horizontally and vertically
                left = Inches(0.5 + (max_width - width) / 2)
                top = Inches(1.5 + (max_height - height) / 2)
            
            # Add image to slide with calculated dimensions
            self.current_slide.shapes.add_picture(
                str(image_path),
                left,
                top,
                width=Inches(width),
                height=Inches(height)
            )
        except Exception as e:
            # If image addition fails, fall back to code display
            # This will be handled by the caller
            raise ConversionError(
                f"Failed to add Mermaid diagram image: {str(e)}",
                format="powerpoint",
                stage="generation"
            ) from e
    
    def _add_blockquote_to_slide(self, node: ASTNode, options: Dict[str, Any]):
        """Add blockquote to current slide."""
        if not self.current_slide:
            self._create_content_slide(None, options)
        
        content = node.content
        
        # Find content placeholder
        content_placeholder = None
        for shape in self.current_slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            text_frame = content_placeholder.text_frame
            p = text_frame.add_paragraph()
            p.left_indent = Inches(0.5)
            # Parse and apply inline formatting (blockquotes are italic by default)
            self._add_formatted_text(p, content, base_font_size=Pt(16), default_italic=True)
    
    def _add_image_to_slide(self, node: ASTNode, options: Dict[str, Any]):
        """Add image to current slide. If text exists, create two-column layout."""
        if not self.current_slide:
            self._create_content_slide(None, options)
        
        image_src = node.attributes.get('src')
        if not image_src:
            return
        
        # Resolve image path
        image_path = Path(image_src)
        if not image_path.is_absolute():
            base_path = options.get('base_path', '.')
            image_path = Path(base_path) / image_path
        
        if not image_path.exists():
            if options.get('skip_missing_images', False):
                return
            raise FileError(
                f"Image file not found: {image_src}",
                file_path=str(image_path),
                operation="read"
            )
        
        try:
            # Check if there's existing text content on the slide
            has_text_content = self._slide_has_text_content()
            
            if has_text_content:
                # Create two-column layout: text on left, image on right
                # Left column: text (0.5" to 4.5" width)
                # Right column: image (5.5" to 9.5" width)
                # Adjust existing content placeholder to left column
                self._reorganize_slide_for_two_columns()
                
                # Calculate dimensions preserving aspect ratio for right column
                max_width = 4.0  # Right column width in inches
                max_height = 4.5  # Available height in inches
                width, height = self._calculate_image_dimensions(image_path, max_width, max_height)
                
                # Position in right column (centered vertically)
                left = Inches(5.5)
                top = Inches(2 + (max_height - height) / 2)  # Center vertically
                
                self.current_slide.shapes.add_picture(
                    str(image_path), 
                    left, 
                    top, 
                    width=Inches(width),
                    height=Inches(height)
                )
            else:
                # No text content, add image centered or full width
                max_width = 8.0  # Full width minus margins
                max_height = 5.5  # Available height
                width, height = self._calculate_image_dimensions(image_path, max_width, max_height)
                
                # Center horizontally and vertically
                left = Inches(1 + (max_width - width) / 2)
                top = Inches(2 + (max_height - height) / 2)
                
                self.current_slide.shapes.add_picture(
                    str(image_path), 
                    left, 
                    top, 
                    width=Inches(width),
                    height=Inches(height)
                )
        
        except Exception as e:
            raise FileError(
                f"Failed to add image: {str(e)}",
                file_path=str(image_path),
                operation="read"
            ) from e
    
    def _slide_has_text_content(self) -> bool:
        """Check if current slide has text content in the content placeholder."""
        if not self.current_slide:
            return False
        
        # Find content placeholder
        content_placeholder = None
        for shape in self.current_slide.shapes:
            try:
                if hasattr(shape, 'placeholder_format'):
                    if shape.placeholder_format.idx == 1:
                        content_placeholder = shape
                        break
            except:
                # Shape doesn't have placeholder_format or it's not accessible
                continue
        
        if not content_placeholder:
            return False
        
        # Check if placeholder has text content
        if hasattr(content_placeholder, 'text_frame'):
            text_frame = content_placeholder.text_frame
            # Check if there are paragraphs with text
            for paragraph in text_frame.paragraphs:
                if paragraph.text and paragraph.text.strip():
                    return True
        
        return False
    
    def _reorganize_slide_for_two_columns(self):
        """Reorganize current slide to use two-column layout with text on left."""
        if not self.current_slide:
            return
        
        # Find content placeholder
        content_placeholder = None
        for shape in self.current_slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if not content_placeholder:
            return
        
        # Adjust content placeholder to left column
        # Left column: 0.5" left margin, 4" width
        # This will cause text to reflow automatically within the new bounds
        content_placeholder.left = Inches(0.5)
        content_placeholder.top = Inches(2)
        content_placeholder.width = Inches(4)
        content_placeholder.height = Inches(5)
        
        # Enable word wrap if not already enabled
        if hasattr(content_placeholder, 'text_frame'):
            content_placeholder.text_frame.word_wrap = True
    
    def _calculate_image_dimensions(self, image_path: Path, max_width: float, max_height: float) -> Tuple[float, float]:
        """
        Calculate image dimensions preserving aspect ratio.
        
        Args:
            image_path: Path to image file
            max_width: Maximum width in inches
            max_height: Maximum height in inches
            
        Returns:
            Tuple of (width, height) in inches
        """
        if not PIL_AVAILABLE:
            # Fallback: return max dimensions if PIL not available
            return (max_width, max_height)
        
        try:
            with PILImage.open(image_path) as img:
                # Get image dimensions in pixels
                img_width, img_height = img.size
                
                # Calculate aspect ratio
                aspect_ratio = img_width / img_height
                
                # Calculate dimensions that fit within max_width and max_height
                # while preserving aspect ratio
                width = max_width
                height = width / aspect_ratio
                
                # If height exceeds max_height, scale based on height instead
                if height > max_height:
                    height = max_height
                    width = height * aspect_ratio
                
                return (width, height)
        except Exception:
            # If image reading fails, return max dimensions
            return (max_width, max_height)
    
    def _add_formatted_text(self, paragraph, text: str, base_font_size=None, default_bold=False, default_italic=False):
        """
        Add text to paragraph with inline markdown formatting applied.
        
        Args:
            paragraph: PowerPoint paragraph object
            text: Text with inline markdown (e.g., **bold**, *italic*)
            base_font_size: Base font size in points
            default_bold: Apply bold by default to all text
            default_italic: Apply italic by default to all text
        """
        # Clear existing runs (new paragraphs have one empty run by default)
        paragraph.clear()
        
        # Parse text into segments with formatting
        segments = self._parse_inline_markdown(text)
        
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment['text']
            
            # Apply base font
            if base_font_size:
                run.font.size = base_font_size
            
            # Apply formatting
            run.font.bold = segment.get('bold', False) or default_bold
            run.font.italic = segment.get('italic', False) or default_italic
            
            if segment.get('code'):
                run.font.name = 'Courier New'
                if base_font_size:
                    run.font.size = Pt(int(base_font_size.pt * 0.9))  # Slightly smaller for code
                run.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)  # Red for code
            
            # Handle links (python-pptx hyperlinks require special handling)
            if segment.get('link_url'):
                run.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)  # Blue for links
                run.font.underline = True
                # Note: Full hyperlink support requires relationship management
                # For now, we just format as a link visually
    
    def _parse_inline_markdown(self, text: str) -> List[Dict]:
        """
        Parse inline markdown into segments with formatting info.
        
        Returns:
            List of segments with text and formatting properties
        """
        segments = []
        
        # Find all formatting markers
        markers = []
        
        # Find bold markers
        for match in self.bold_pattern.finditer(text):
            markers.append({
                'start': match.start(),
                'end': match.end(),
                'type': 'bold',
                'content': match.group(1) or match.group(2)
            })
        
        # Find italic markers (not overlapping with bold)
        for match in self.italic_pattern.finditer(text):
            # Check if this is part of a bold marker
            is_bold_part = False
            for bold_marker in markers:
                if bold_marker['start'] <= match.start() < bold_marker['end']:
                    is_bold_part = True
                    break
            if not is_bold_part:
                markers.append({
                    'start': match.start(),
                    'end': match.end(),
                    'type': 'italic',
                    'content': match.group(1) or match.group(2)
                })
        
        # Find code markers
        for match in self.code_pattern.finditer(text):
            markers.append({
                'start': match.start(),
                'end': match.end(),
                'type': 'code',
                'content': match.group(1)
            })
        
        # Find link markers
        for match in self.link_pattern.finditer(text):
            markers.append({
                'start': match.start(),
                'end': match.end(),
                'type': 'link',
                'content': match.group(1),
                'url': match.group(2)
            })
        
        # Sort markers by position
        markers.sort(key=lambda x: x['start'])
        
        # Build segments
        pos = 0
        for marker in markers:
            # Add text before marker
            if marker['start'] > pos:
                segments.append({
                    'text': text[pos:marker['start']],
                    'bold': False,
                    'italic': False,
                    'code': False
                })
            
            # Add formatted segment
            segment = {
                'text': marker['content'],
                'bold': marker['type'] == 'bold',
                'italic': marker['type'] == 'italic',
                'code': marker['type'] == 'code'
            }
            
            if marker['type'] == 'link':
                segment['link_url'] = marker['url']
            
            segments.append(segment)
            pos = marker['end']
        
        # Add remaining text
        if pos < len(text):
            segments.append({
                'text': text[pos:],
                'bold': False,
                'italic': False,
                'code': False
            })
        
        # If no markers found, return plain text
        if not segments:
            segments.append({
                'text': text,
                'bold': False,
                'italic': False,
                'code': False
            })
        
        return segments
    
    def _find_first_heading(self, node: ASTNode, level: int) -> Optional[ASTNode]:
        """Find first heading of specified level."""
        if node.node_type == NodeType.HEADING and node.level == level:
            return node
        
        for child in node.children:
            result = self._find_first_heading(child, level)
            if result:
                return result
        
        return None

