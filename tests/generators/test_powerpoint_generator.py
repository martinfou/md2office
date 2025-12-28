"""
Tests for PowerPoint Document Generator

Implements Story 3.5: PowerPoint Testing and Validation
"""

import pytest
import sys
from pathlib import Path

# Add src to path (now in subdirectory, go up two levels)
sys.path.insert(0, str(Path(__file__).parent.parent.parent / 'src'))

from md2office.parser import MarkdownParser, ASTBuilder
from md2office.generators import PowerPointGenerator
from md2office.errors import ConversionError


class TestPowerPointGenerator:
    """Test suite for PowerPoint generator."""
    
    @pytest.fixture
    def powerpoint_generator(self):
        """Create PowerPoint generator instance."""
        try:
            return PowerPointGenerator()
        except ImportError:
            pytest.skip("python-pptx not available")
    
    @pytest.fixture
    def sample_markdown(self):
        """Sample markdown content for testing."""
        return """# Test Presentation

This is a test presentation.

## Section 1

Some content here.

### Subsection

- Item 1
- Item 2

## Section 2

```python
def hello():
    print("Hello, World!")
```

| Header 1 | Header 2 |
|----------|----------|
| Cell 1   | Cell 2   |
"""
    
    def test_generator_initialization(self, powerpoint_generator):
        """Test PowerPoint generator initialization."""
        assert powerpoint_generator is not None
    
    def test_generate_basic_presentation(self, powerpoint_generator, sample_markdown):
        """Test generating a basic PowerPoint presentation."""
        # Parse markdown
        parser = MarkdownParser()
        tokens = parser.parse(sample_markdown)
        
        # Build AST
        builder = ASTBuilder()
        ast = builder.build(tokens)
        
        # Generate PowerPoint presentation
        options = {}
        pptx_bytes = powerpoint_generator.generate(ast, options)
        
        # Verify presentation was generated
        assert pptx_bytes is not None
        assert len(pptx_bytes) > 0
        
        # Verify it's a valid PPTX file (starts with ZIP signature)
        assert pptx_bytes[:2] == b'PK'  # ZIP file signature
    
    def test_generate_with_metadata(self, powerpoint_generator):
        """Test generating presentation with metadata."""
        markdown = """---
title: Test Presentation
author: Test Author
---

# Test Presentation

Content here.
"""
        parser = MarkdownParser()
        tokens = parser.parse(markdown)
        builder = ASTBuilder()
        ast = builder.build(tokens)
        
        options = {}
        pptx_bytes = powerpoint_generator.generate(ast, options)
        assert pptx_bytes is not None
    
    def test_generate_with_style_preset(self, powerpoint_generator, sample_markdown):
        """Test generating presentation with style preset."""
        parser = MarkdownParser()
        tokens = parser.parse(sample_markdown)
        builder = ASTBuilder()
        ast = builder.build(tokens)
        
        options = {'style': 'professional'}
        pptx_bytes = powerpoint_generator.generate(ast, options)
        assert pptx_bytes is not None
    
    def test_slide_structure(self, powerpoint_generator, sample_markdown):
        """Test that slides are created correctly."""
        parser = MarkdownParser()
        tokens = parser.parse(sample_markdown)
        builder = ASTBuilder()
        ast = builder.build(tokens)
        
        options = {}
        pptx_bytes = powerpoint_generator.generate(ast, options)
        
        # Verify presentation structure
        from pptx import Presentation
        from io import BytesIO
        
        pres = Presentation(BytesIO(pptx_bytes))
        assert len(pres.slides) > 0  # Should have at least title slide
    
    def test_file_extension(self, powerpoint_generator):
        """Test file extension method."""
        assert powerpoint_generator.get_file_extension() == ".pptx"


class TestPowerPointGeneratorIntegration:
    """Integration tests for PowerPoint generator."""
    
    @pytest.fixture
    def powerpoint_generator(self):
        """Create PowerPoint generator instance."""
        try:
            return PowerPointGenerator()
        except ImportError:
            pytest.skip("python-pptx not available")
    
    def test_full_conversion_pipeline(self, powerpoint_generator):
        """Test full conversion pipeline from markdown to PowerPoint."""
        try:
            from md2office.router import ConversionPipeline
            from md2office.generators import PowerPointGenerator
            
            pipeline = ConversionPipeline()
            pptx_gen = PowerPointGenerator()
            pipeline.register_generator('powerpoint', pptx_gen)
            
            markdown = """# Test Presentation

This is a **test** presentation.
"""
            results = pipeline.convert(markdown, ['powerpoint'])
            
            assert 'powerpoint' in results
            assert results['powerpoint'] is not None
            assert len(results['powerpoint']) > 0
        
        except ImportError:
            pytest.skip("python-pptx not available")
    
    def test_no_code_overlap_with_text(self):
        """Test that code blocks never overlap with bullet points or text content."""
        try:
            from md2office.generators import PowerPointGenerator
            from pptx import Presentation
            from pptx.util import Inches
            from io import BytesIO
            
            generator = PowerPointGenerator()
            
            # Create markdown with various scenarios: text before code, code before text
            markdown = """# Test Presentation

## Text Before Code

- First bullet point
- Second bullet point
- Third bullet point

```python
def example_function():
    return "code here"
```

## Code Before Text

```python
def another_function():
    return "code first"
```

- Bullet point after code
- Another bullet point

## Mixed Content

Some paragraph text here.

```python
code_block()
```

- More bullets
- Even more bullets

## Image with Text

- Point 1
- Point 2

![test](image.png)

## Text with Image

![test](image.png)

- Bullet after image
"""
            
            parser = MarkdownParser()
            tokens = parser.parse(markdown)
            builder = ASTBuilder()
            ast = builder.build(tokens)
            
            options = {'skip_missing_images': True}  # Skip missing images for test
            pptx_bytes = generator.generate(ast, options)
            
            # Load presentation
            pres = Presentation(BytesIO(pptx_bytes))
            
            # Check each slide for overlaps
            overlap_issues = []
            
            for slide_idx, slide in enumerate(pres.slides):
                # Find content placeholder (text content)
                content_placeholder = None
                code_blocks = []
                images = []
                
                for shape in slide.shapes:
                    try:
                        # Check if it's the content placeholder (idx 1)
                        if hasattr(shape, 'placeholder_format'):
                            try:
                                if shape.placeholder_format.idx == 1:
                                    content_placeholder = shape
                            except:
                                pass
                        
                        # Check if it's a code block (textbox with gray background)
                        if hasattr(shape, 'fill') and hasattr(shape, 'text_frame'):
                            try:
                                if hasattr(shape.fill, 'fore_color'):
                                    if hasattr(shape.fill.fore_color, 'rgb'):
                                        rgb = shape.fill.fore_color.rgb
                                        if rgb and len(rgb) >= 3:
                                            # Code blocks use light gray background (#F5F5F5)
                                            if rgb[0] >= 240 and rgb[1] >= 240 and rgb[2] >= 240:
                                                code_blocks.append(shape)
                            except:
                                pass
                        
                        # Check if it's an image
                        if hasattr(shape, 'image'):
                            images.append(shape)
                    
                    except:
                        continue
                
                # Check if slide has both text content and code/images
                has_text = False
                if content_placeholder and hasattr(content_placeholder, 'text_frame'):
                    text_frame = content_placeholder.text_frame
                    for paragraph in text_frame.paragraphs:
                        if paragraph.text and paragraph.text.strip():
                            has_text = True
                            break
                
                has_code_or_image = len(code_blocks) > 0 or len(images) > 0
                
                if has_text and has_code_or_image:
                    # Verify two-column layout
                    # Text should be in left column (content placeholder x < 5 inches)
                    # Code/images should be in right column (x >= 5 inches)
                    
                    # Check content placeholder position (should be left column)
                    content_right_edge = content_placeholder.left + content_placeholder.width
                    if content_right_edge > Inches(5):
                        overlap_issues.append(
                            f"Slide {slide_idx + 1}: Content placeholder extends into right column "
                            f"(right edge at {content_right_edge.inches:.2f} inches)"
                        )
                    
                    # Check code blocks position (should be right column)
                    for code_idx, code_block in enumerate(code_blocks):
                        if code_block.left < Inches(5):
                            overlap_issues.append(
                                f"Slide {slide_idx + 1}: Code block {code_idx + 1} is in left column "
                                f"(left edge at {code_block.left.inches:.2f} inches, should be >= 5 inches)"
                            )
                        
                        # Check if code block overlaps with content placeholder
                        code_left = code_block.left
                        code_right = code_block.left + code_block.width
                        content_left = content_placeholder.left
                        content_right_edge = content_placeholder.left + content_placeholder.width
                        
                        # Check for horizontal overlap
                        if not (code_right <= content_left or code_left >= content_right_edge):
                            overlap_issues.append(
                                f"Slide {slide_idx + 1}: Code block {code_idx + 1} overlaps with text "
                                f"(code: {code_left.inches:.2f}-{code_right.inches:.2f}in, "
                                f"text: {content_left.inches:.2f}-{content_right_edge.inches:.2f}in)"
                            )
                    
                    # Check images position (should be right column)
                    for img_idx, image in enumerate(images):
                        if image.left < Inches(5):
                            overlap_issues.append(
                                f"Slide {slide_idx + 1}: Image {img_idx + 1} is in left column "
                                f"(left edge at {image.left.inches:.2f} inches, should be >= 5 inches)"
                            )
                        
                        # Check if image overlaps with content placeholder
                        img_left = image.left
                        img_right = image.left + image.width
                        content_left = content_placeholder.left
                        content_right_edge = content_placeholder.left + content_placeholder.width
                        
                        # Check for horizontal overlap
                        if not (img_right <= content_left or img_left >= content_right_edge):
                            overlap_issues.append(
                                f"Slide {slide_idx + 1}: Image {img_idx + 1} overlaps with text "
                                f"(image: {img_left.inches:.2f}-{img_right.inches:.2f}in, "
                                f"text: {content_left.inches:.2f}-{content_right_edge.inches:.2f}in)"
                            )
            
            # Assert no overlap issues found
            if overlap_issues:
                error_msg = "Found code/image overlap issues:\n" + "\n".join(f"  - {issue}" for issue in overlap_issues)
                pytest.fail(error_msg)
        
        except ImportError:
            pytest.skip("python-pptx not available")
    
    def test_no_code_overlap_in_presentation_file(self):
        """Test that the actual presentation file has no code/image overlaps with text."""
        try:
            from md2office.generators import PowerPointGenerator
            from pptx import Presentation
            from pptx.util import Inches
            from io import BytesIO
            
            # Use the actual presentation file from examples
            test_file = Path(__file__).parent.parent.parent / 'examples' / 'input' / 'presentation-session-1-quick-wins.md'
            
            if not test_file.exists():
                pytest.skip(f"Test file not found: {test_file}")
            
            # Generate PowerPoint from the actual file
            generator = PowerPointGenerator()
            
            markdown_content = test_file.read_text(encoding='utf-8')
            parser = MarkdownParser()
            tokens = parser.parse(markdown_content)
            builder = ASTBuilder()
            ast = builder.build(tokens)
            
            options = {'skip_missing_images': True, 'base_path': str(test_file.parent)}
            pptx_bytes = generator.generate(ast, options)
            
            # Load presentation
            pres = Presentation(BytesIO(pptx_bytes))
            
            # Check each slide for overlaps
            overlap_issues = []
            
            for slide_idx, slide in enumerate(pres.slides):
                # Find content placeholder (text content)
                content_placeholder = None
                code_blocks = []
                images = []
                
                for shape in slide.shapes:
                    try:
                        # Check if it's the content placeholder (idx 1)
                        if hasattr(shape, 'placeholder_format'):
                            try:
                                if shape.placeholder_format.idx == 1:
                                    content_placeholder = shape
                            except:
                                pass
                        
                        # Check if it's a code block (textbox with gray background)
                        if hasattr(shape, 'fill') and hasattr(shape, 'text_frame'):
                            try:
                                if hasattr(shape.fill, 'fore_color'):
                                    if hasattr(shape.fill.fore_color, 'rgb'):
                                        rgb = shape.fill.fore_color.rgb
                                        if rgb and len(rgb) >= 3:
                                            # Code blocks use light gray background (#F5F5F5)
                                            if rgb[0] >= 240 and rgb[1] >= 240 and rgb[2] >= 240:
                                                code_blocks.append(shape)
                            except:
                                pass
                        
                        # Check if it's an image
                        if hasattr(shape, 'image'):
                            images.append(shape)
                    
                    except:
                        continue
                
                # Check if slide has both text content and code/images
                has_text = False
                if content_placeholder and hasattr(content_placeholder, 'text_frame'):
                    text_frame = content_placeholder.text_frame
                    for paragraph in text_frame.paragraphs:
                        if paragraph.text and paragraph.text.strip():
                            has_text = True
                            break
                
                has_code_or_image = len(code_blocks) > 0 or len(images) > 0
                
                if has_text and has_code_or_image:
                    # Verify two-column layout
                    # Text should be in left column (content placeholder x < 5 inches)
                    # Code/images should be in right column (x >= 5 inches)
                    
                    # Check content placeholder position (should be left column)
                    content_right_edge = content_placeholder.left + content_placeholder.width
                    if content_right_edge > Inches(5):
                        overlap_issues.append(
                            f"Slide {slide_idx + 1}: Content placeholder extends into right column "
                            f"(right edge at {content_right_edge.inches:.2f} inches)"
                        )
                    
                    # Check code blocks position (should be right column)
                    for code_idx, code_block in enumerate(code_blocks):
                        if code_block.left < Inches(5):
                            overlap_issues.append(
                                f"Slide {slide_idx + 1}: Code block {code_idx + 1} is in left column "
                                f"(left edge at {code_block.left.inches:.2f} inches, should be >= 5 inches)"
                            )
                        
                        # Check if code block overlaps with content placeholder
                        code_left = code_block.left
                        code_right = code_block.left + code_block.width
                        content_left = content_placeholder.left
                        content_right_edge = content_placeholder.left + content_placeholder.width
                        
                        # Check for horizontal overlap
                        if not (code_right <= content_left or code_left >= content_right_edge):
                            overlap_issues.append(
                                f"Slide {slide_idx + 1}: Code block {code_idx + 1} overlaps with text "
                                f"(code: {code_left.inches:.2f}-{code_right.inches:.2f}in, "
                                f"text: {content_left.inches:.2f}-{content_right_edge.inches:.2f}in)"
                            )
                    
                    # Check images position (should be right column)
                    for img_idx, image in enumerate(images):
                        if image.left < Inches(5):
                            overlap_issues.append(
                                f"Slide {slide_idx + 1}: Image {img_idx + 1} is in left column "
                                f"(left edge at {image.left.inches:.2f} inches, should be >= 5 inches)"
                            )
                        
                        # Check if image overlaps with content placeholder
                        img_left = image.left
                        img_right = image.left + image.width
                        content_left = content_placeholder.left
                        content_right_edge = content_placeholder.left + content_placeholder.width
                        
                        # Check for horizontal overlap
                        if not (img_right <= content_left or img_left >= content_right_edge):
                            overlap_issues.append(
                                f"Slide {slide_idx + 1}: Image {img_idx + 1} overlaps with text "
                                f"(image: {img_left.inches:.2f}-{img_right.inches:.2f}in, "
                                f"text: {content_left.inches:.2f}-{content_right_edge.inches:.2f}in)"
                            )
            
            # Assert no overlap issues found
            if overlap_issues:
                error_msg = f"Found {len(overlap_issues)} code/image overlap issue(s) in presentation:\n" + "\n".join(f"  - {issue}" for issue in overlap_issues)
                pytest.fail(error_msg)
        
        except ImportError:
            pytest.skip("python-pptx not available")

